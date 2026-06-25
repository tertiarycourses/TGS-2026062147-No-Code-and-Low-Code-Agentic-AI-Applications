#!/usr/bin/env python3
"""Highly professional 3-day course deck (150-200 slides): courseware/n8n-slides.pptx
Admin front matter + full n8n key concepts (mined from the reference deck) +
per-activity overview, workflow screenshot and step-by-step slides + breaks."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import os as _os
REPO = _os.path.dirname(_os.path.dirname(_os.path.dirname(_os.path.dirname(_os.path.abspath(__file__)))))

NAVY=RGBColor(0x0B,0x12,0x20); BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81)
AMBER=RGBColor(0xF5,0x9E,0x0B); INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72)
LIGHT=RGBColor(0xF5,0xF8,0xFC); WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0)
VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]; FONT="Arial"
def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=color
    sh.line.fill.background() if line is None else (setattr(sh.line.color,'rgb',line),setattr(sh.line,'width',Pt(1)))
    sh.shadow.inherit=False; return sh
def oval(s,x,y,w,h,color):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,x,y,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=color
    sh.line.fill.background(); sh.shadow.inherit=False; return sh
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(0); tf.margin_top=tf.margin_bottom=Pt(0)
    if runs and not isinstance(runs[0],list): runs=[runs]
    first=True
    for para in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.space_after=Pt(space); p.space_before=Pt(0)
        for (t,sz,col,bold) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.margin_left=tf.margin_right=Pt(0)
    first=True
    for it in items:
        lvl=0; text=it
        if isinstance(it,tuple): text,lvl=it
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_after=Pt(gap); p.space_before=Pt(0)
        rm=p.add_run(); rm.text=("   "*lvl)+("•  " if lvl==0 else "–  ")
        rm.font.size=Pt(size); rm.font.color.rgb=(mcolor if lvl==0 else GREY); rm.font.bold=True; rm.font.name=FONT
        r=p.add_run(); r.text=text; r.font.size=Pt(size-lvl); r.font.color.rgb=color; r.font.name=FONT
    return tb
PAGE={"n":0}
def footer(s,dark=False):
    PAGE["n"]+=1; c=WHITE if dark else GREY
    txt(s,Inches(0.55),Inches(7.08),Inches(6.2),Inches(0.3),[[("Agentic AI Automation with n8n  ·  TGS-2023035977",8.5,c,False)]])
    txt(s,Inches(6.0),Inches(7.08),Inches(6),Inches(0.3),[[("© 2026 Tertiary Infotech Academy Pte Ltd",8.5,c,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(11.9),Inches(7.08),Inches(0.9),Inches(0.3),[[(str(PAGE["n"]),9,c,True)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,Inches(0.55),Inches(0.60),Inches(0.14),Inches(0.66),kcolor)
    y=Inches(0.55)
    if kicker:
        txt(s,Inches(0.8),Inches(0.5),Inches(11.8),Inches(0.34),[[(kicker,14,kcolor,True)]]); y=Inches(0.82)
    txt(s,Inches(0.8),y,Inches(12.0),Inches(0.85),[[(title,29,INK,True)]])
    rect(s,Inches(0.8),Inches(1.55),Inches(12.0),Pt(2),LINE)

ASSETS=f"{REPO}/courseware/assets"
def cover():
    import os
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE)
    rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    # Tertiary Infotech Academy logo (top-left) + n8n course logo (top-right)
    if os.path.exists(f"{ASSETS}/tertiary-infotech-logo.png"):
        s.shapes.add_picture(f"{ASSETS}/tertiary-infotech-logo.png",Inches(0.85),Inches(0.7),height=Inches(1.05))
    if os.path.exists(f"{ASSETS}/n8n-course-logo.png"):
        s.shapes.add_picture(f"{ASSETS}/n8n-course-logo.png",Inches(11.4),Inches(0.72),height=Inches(1.0))
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("LEARNER GUIDE  ·  COURSE SLIDES",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.9),Inches(12.0),Inches(1.6),[[("Agentic AI Automation with n8n",46,INK,True)]])
    rect(s,Inches(0.92),Inches(4.3),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(4.65),Inches(12),Inches(1.4),
        [[("WSQ Course Code: TGS-2023035977",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.55),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])
def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.5),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.3),Inches(0.7),Inches(2.5),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=slide(); head(s,title,kicker); bullets(s,Inches(0.85),Inches(1.9),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=slide(); head(s,title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=17)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=17,mcolor=TEAL); footer(s); return s
def img_slide(title,img,caption,kicker="WORKFLOW"):
    s=slide(); head(s,title,kicker,TEAL)
    pic=s.shapes.add_picture(img,Inches(1.0),Inches(1.85),width=Inches(11.3))
    maxh=Inches(4.55)
    if pic.height>maxh:
        s.shapes._spTree.remove(pic._element)
        pic=s.shapes.add_picture(img,Inches(1.0),Inches(1.85),height=maxh)
        pic.left=int((SW-pic.width)/2)
    else:
        pic.left=int((SW-pic.width)/2)
    txt(s,Inches(0.8),Inches(6.55),Inches(11.7),Inches(0.4),[[(caption,13,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def website_slide(title,img,items,kicker,note=""):
    import os
    s=slide(); head(s,title,kicker,BLUE)
    bullets(s,Inches(0.85),Inches(1.95),Inches(5.4),Inches(4.6),items,size=18)
    if os.path.exists(img):
        rect(s,Inches(6.42),Inches(1.92),Inches(6.18),Inches(3.92),LINE)
        pic=s.shapes.add_picture(img,Inches(6.5),Inches(2.0),width=Inches(6.0))
        if pic.height>Inches(3.78):
            s.shapes._spTree.remove(pic._element); pic=s.shapes.add_picture(img,Inches(6.5),Inches(2.0),height=Inches(3.78)); pic.left=int(Inches(6.5)+(Inches(6.0)-pic.width)/2)
    footer(s); return s
def gallery_slide(title,imgs,captions,kicker):
    import os
    s=slide(); head(s,title,kicker,TEAL)
    imgs2=[(i,c) for i,c in zip(imgs,captions) if os.path.exists(i)]; n=len(imgs2)
    if not n: footer(s); return s
    slot=11.6/n; cell=Inches(slot)
    for k,(im,cap) in enumerate(imgs2):
        x0=Inches(0.85)+Inches(slot*k)
        pic=s.shapes.add_picture(im,x0+Inches(0.1),Inches(2.0),width=cell-Inches(0.3))
        if pic.height>Inches(3.7):
            s.shapes._spTree.remove(pic._element); pic=s.shapes.add_picture(im,x0,Inches(2.0),height=Inches(3.7))
        pic.left=int(x0+(cell-pic.width)/2)
        txt(s,x0,Inches(5.95),cell,Inches(0.5),[[(cap,13,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,40,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.7),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
def cards3(title,cards,kicker):
    s=slide(); head(s,title,kicker)
    xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT)
        rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def activity_overview(tag,title,desc,build,nodes,kicker):
    s=slide(); head(s,title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.82),Inches(1.7),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.87),Inches(1.7),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,21,INK,False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),LIGHT)
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll build",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.6),[[(build,18,INK,True)]])
    txt(s,Inches(1.1),Inches(5.6),Inches(11),Inches(0.5),[[("Key nodes:  ",13,GREY,True),(nodes,13,GREY,False)]]); footer(s); return s
def step_slide(kicker,act_title,n,total,text):
    s=slide(); head(s,act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.6),Inches(1.5),Inches(1.5),TEAL)
    txt(s,Inches(0.85),Inches(2.86),Inches(1.5),Inches(1.0),[[(str(n),40,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.7),Inches(2.6),Inches(10.0),Inches(1.5),[[(text,24,INK,False)]],anchor=MSO_ANCHOR.MIDDLE); footer(s); return s
def test_slide(act_title,text,kicker):
    s=slide(); head(s,act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.4),[[(text,18,INK,False)]]); footer(s); return s
def brk(kind,dur,color):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

IMG=lambda p: f"{REPO}/{p}"

# ============================================================ BUILD
cover()

# ---------- ADMIN ----------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator will display the digital attendance QR code generated from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
content("About the Trainer",[
 "Dr. Alfred Ang — Principal Trainer, Tertiary Infotech Academy Pte. Ltd.",
 "PhD; specialises in Artificial Intelligence, automation and software engineering.",
 "Designs and delivers WSQ courses on AI agents, automation (n8n) and app development.",
 "Founder and lead instructor at Tertiary Infotech / Tertiary Courses."],kicker="YOUR TRAINER")
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with automation or AI tools (if any).",
 "What you want to automate after this course."],kicker="ICE-BREAKER")
content("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","Step out quietly for calls or toilet breaks.",
 "75% attendance is required."])
content("LMS / TMS",[
 "Access your course materials, attendance and assessment on the LMS/TMS portal.",
 "Portal: https://lms-tms.tertiaryinfotech.com",
 "Download the slides and Learner Guide for reference during the open-book assessment."],kicker="COURSE PORTAL")
two_col("Lesson Plan — 3 Days, 8 hours/day",[
 ("Day 1 — Workflow Automation + AI Agents",0),("Topic 1: n8n basics + Activities 1, 2, 3a, 3b",1),
 ("Topic 2: AI Agents — Activities 4a, 4b",1),("Day 2 — RAG, Webhooks, APIs",0),
 ("RAG (Activity 5)",1),("Webhooks (Activity 6) · APIs (Activity 7)",1)],
 [("Day 3 — Security & Capstone",0),("Human-in-the-loop + Guardrails (Activity 8)",1),
 ("Mini Capstone + presentations",1),("Daily timing",0),
 ("9:30am–6:30pm · 1-hour lunch",1),("Short tea breaks within each day",1)],
 kicker="SCHEDULE",lhead="Days 1–2",rhead="Day 3 & timing")
content("Learning Outcomes",[
 "LO1: Build automated workflows in n8n using triggers, actions, nodes and flows.",
 "LO2: Design AI agents with LLMs, memory and tools, triggered from chat / Telegram.",
 "LO3: Apply Retrieval-Augmented Generation (RAG) to ground agents in your documents.",
 "LO4: Integrate external systems via webhooks, APIs and HTTP requests.",
 "LO5: Apply human-in-the-loop and guardrails to secure agentic automations."],kicker="WHAT YOU'LL ACHIEVE")
content("Assessment",[
 "Written Assessment (SAQ) — 1 hour.","Practical Performance (PP) — 1 hour.",
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 "A mini capstone project is presented on Day 3.","An appeal process is available if required."],kicker="FINAL ASSESSMENT")
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."])

# ---------- TOPIC 1: WORKFLOW AUTOMATION ----------
section("TOPIC 1","Workflow Automation with n8n","01","Triggers · Actions · Nodes · Flows")
content("What is n8n?",[
 "A workflow automation platform — connect apps and services with little or no code.",
 "Visual editor: drag nodes onto a canvas and wire them together.",
 "Runs in the cloud (trial) or self-hosted locally with Docker.",
 "400+ integrations, HTTP/code nodes, and AI/LangChain nodes for agents."],kicker="OVERVIEW")
cards3("Features of n8n",[
 (BLUE,"Build visually",["400+ app integrations","Drag-and-drop node canvas","No-code to full-code (JS/Python)"]),
 (TEAL,"Run anywhere",["Cloud or self-host (Docker)","Webhooks, schedules, queues","Version & export workflows"]),
 (VIOLET,"AI-native",["AI Agent + LangChain nodes","Memory, tools, RAG, MCP","Use any LLM (OpenAI, Gemini...)"])],kicker="FEATURES")
big_statement("One canvas to connect apps, data and AI.","n8n lets you automate work and build AI agents without gluing code together by hand.","WHY USE n8n",color=BLUE)
website_slide("n8n.io",IMG("courseware/assets/site-n8n.png"),
 ["n8n is an open, fair-code workflow automation platform.","Free cloud trial at n8n.io, or self-host with Docker.","Same visual editor in both."],
 kicker="THE PLATFORM")
content("Why Automate?",[
 "Remove repetitive manual work and human error.",
 "Connect systems that don't normally talk to each other.",
 "React to events in real time (forms, webhooks, schedules).",
 "Add AI to make workflows understand language and make decisions."],kicker="MOTIVATION")
content("Setting Up n8n",[
 "Option A — Cloud trial: sign up at n8n.io and start in the editor immediately.",
 "Option B — Local Docker Compose (persistent): docker compose up -d → http://localhost:5678.",
 "Trial Data Tables are not permanent — store anything you keep externally (Google Sheets).",
 "See labs/n8n-installation/docker-compose.yml in the course repo."],kicker="GET READY")
content("Credential Setup",[
 "Add credentials once under Credentials → Add credential.",
 "Gmail / Outlook (OAuth2) for email; OpenAI / Gemini for AI.",
 "Telegram bot token (from @BotFather); Google Sheets (OAuth2).",
 "Imported workflows reference credential names — re-select your own after importing."],kicker="ACCOUNTS & KEYS")
content("The n8n Editor",[
 "Canvas — where you place and connect nodes.",
 "Node panel — search 400+ nodes by name.",
 "Execute Workflow — run and inspect data at each step.",
 "Each node shows input and output data as JSON."],kicker="UI TOUR")
two_col("n8n Nodes",[
 ("Trigger nodes — start a workflow",0),("Form, Webhook, Telegram, Schedule, Manual",1),
 ("Action nodes — do something",0),("Gmail, Google Sheets, HTTP Request, Data Table",1)],
 [("Logic nodes — control the flow",0),("IF, Switch, Merge, Split Out, Code, Edit Fields",1),
 ("AI nodes — reason and act",0),("AI Agent, Chat Model, Memory, Vector Store, Tools",1)],
 kicker="THE BUILDING BLOCKS",lhead="Triggers & Actions",rhead="Logic & AI")
two_col("Triggers Available in n8n",[
 ("Manual Trigger - run on demand (testing)",0),
 ("Schedule Trigger - run on a timer / cron",0),
 ("Form Trigger - a hosted web form",0),
 ("Webhook - external apps call a URL",0)],
 [("Chat Trigger - built-in chat UI",0),
 ("Telegram Trigger - messages to your bot",0),
 ("Email (IMAP) - on new email",0),
 ("App triggers - Gmail, Sheets, Notion, ...",0)],
 kicker="WHEN A WORKFLOW RUNS",lhead="Core triggers",rhead="Chat & app triggers")
two_col("Key Nodes in n8n",[
 ("HTTP Request - call any API",0),
 ("Gmail / Outlook - send email",0),
 ("Google Sheets / Data Table - store rows",0),
 ("Code - run JavaScript / Python",0),
 ("Edit Fields (Set) - reshape data",0)],
 [("IF / Switch - branch on conditions",0),
 ("Merge / Split Out - combine / expand",0),
 ("AI Agent - reason with tools + memory",0),
 ("Vector Store - RAG retrieval",0),
 ("Respond to Webhook - reply to caller",0)],
 kicker="THE WORKHORSE NODES",lhead="Actions",rhead="Logic & AI")
content("Triggers and Actions",[
 "A Trigger is the first node — it decides WHEN a workflow runs.",
 "Manual & Schedule triggers for testing and time-based jobs.",
 "Form, Webhook and Telegram triggers respond to external events.",
 "Actions run after the trigger — send email, store data, call an API."],kicker="CORE CONCEPT")
content("Execution Modes",[
 "Test execution — run manually in the editor to inspect data.",
 "Production execution — runs automatically when Active.",
 "Each run is recorded in the Execution History for debugging."],kicker="HOW IT RUNS")
content("Data Structure & JSON",[
 "Data flows between nodes as a list of items, each a JSON object.",
 "Each item has a json field with key/value pairs.",
 "Nodes read the previous node's output and add to it."],kicker="DATA")
content("What is JSON?",[
 "JSON (JavaScript Object Notation) is a simple text format for data.",
 'It stores key/value pairs:  { "name": "Alice", "age": 30 }.',
 "Values can be text, numbers, true/false, lists [...] or nested objects {...}.",
 "n8n passes data between nodes as JSON - and most APIs send/receive JSON.",
 "Read a value with an expression, e.g.  {{ $json.name }}."],kicker="THE DATA FORMAT")
content("Expressions & Data Mapping",[
 "Expressions pull data from earlier nodes: {{ $json.Name }}.",
 "Drag fields from the input panel to map them automatically.",
 "Use expressions in any field — subject lines, URLs, conditions."],kicker="DATA MAPPING")
content("Pin Data & Execution History",[
 "Pin data to freeze a node's output while you build downstream nodes.",
 "Edit output to test different scenarios without re-running triggers.",
 "Execution History shows every run, its data, and any errors."],kicker="DEBUGGING")
content("Transforming Data",[
 "Edit Fields (Set) — add, rename or reshape fields.",
 "Code node — run JavaScript/Python for custom logic.",
 "Split Out — turn one item with a list into many items.",
 "Merge — combine data from two branches (append, combine, choose)."],kicker="SHAPING DATA")
content("Conditional Logic — IF & Switch",[
 "IF node — two outputs: true and false, based on a condition.",
 "Switch node — many outputs for multiple cases.",
 "Use them to route submissions, approvals, or message types."],kicker="DECISIONS")
content("Sub-Workflows",[
 "Break big automations into reusable smaller workflows.",
 "Call a sub-workflow with the Execute Workflow node.",
 "Keeps workflows readable and lets you reuse common logic."],kicker="STRUCTURE")

def activity_block(a):
    activity_overview(a["tag"],a["title"],a["desc"],a["build"],a["nodes"],a["kicker"])
    if a.get("img"): img_slide(a["title"]+" — Workflow",IMG(a["img"]),a.get("cap",""),a["kicker"])
    steps=a.get("steps",[]); total=len(steps)
    for i,st in enumerate(steps,1): step_slide(a["kicker"],a["title"],i,total,st)
    if a.get("test"): test_slide(a["title"],a["test"],a["kicker"])

K1="TOPIC 1 · WORKFLOW AUTOMATION"
activity_block(dict(tag="ACT 1",title="Activity 1 — Flyer with QR Code (Form → Email)",kicker=K1,
 desc="Collect a visitor's details with an n8n Form and email them to an admin, then turn the form URL into a QR code on an event flyer (group activity).",
 build="Form Trigger  →  Gmail (Send)",nodes="formTrigger, gmail",
 img="labs/activity1-flyer-form/Activity1-Flyer-Form.png",cap="Form Trigger → Gmail",
 steps=["Create a new workflow and add an n8n Form Trigger; set a Form Title.",
   "Add four fields: Name, Email, Phone, Message (mark Name & Email required).",
   "Add a Gmail node (Send a Message) after the Form Trigger.",
   "Set To = admin address, Subject = New RSVP from {{ $json.Name }}, and a body using the fields.",
   "Select your Gmail credential, Save, and toggle the workflow Active.",
   "Copy the Form's Production URL; paste it into the QR generator (alfredang.github.io/qrcodegenerator) and put the QR on a flyer.",
   "Group activity: design an event flyer (e.g. bowling night) with the QR code, then present it."],
 test="Scan the QR code, submit the form, and confirm the admin inbox receives the email."))
gallery_slide("Sample Flyers from Past Students",
 [IMG("labs/activity1-flyer-form/flyer-sample1-preview.png"),
  IMG("labs/activity1-flyer-form/flyer-sample2-preview.png"),
  IMG("labs/activity1-flyer-form/flyer-sample3.jpeg")],
 ["Network event","Event poster","Bowling party"],kicker="GROUP ACTIVITY")
activity_block(dict(tag="ACT 2",title="Activity 2 — Capture Submissions in a Data Table",kicker=K1,
 desc="Extend Activity 1 so every submission is also saved into an n8n Data Table — your first taste of storing data, not just forwarding it.",
 build="Form Trigger  →  Gmail  +  Data Table (Insert row)",nodes="formTrigger, gmail, dataTable",
 img="labs/activity2-data-table/Activity2-Data-Table.png",cap="Form → Gmail + Data Table",
 steps=["In n8n open Data Tables and create a table 'RSVPs' with columns Name, Email, Phone, Message.",
   "Continue from your Activity 1 workflow.",
   "Add a Data Table node (Insert Row) connected after the Form Trigger, alongside Gmail.",
   "Map each form field to its column with expressions, e.g. Name → {{ $json.Name }}.",
   "Save and keep the workflow Active."],
 test="Submit the form and confirm a new row appears in the RSVPs Data Table and the email still sends."))
activity_block(dict(tag="ACT 3a",title="Activity 3a — Conditional Response (Data Table)",kicker=K1,
 desc="Add decision-making with an IF node. If the visitor is attending, save the date to the Data Table; if not, send a polite thank-you email.",
 build="Form Trigger  →  IF  →  Data Table / Gmail",nodes="formTrigger, if, dataTable, gmail",
 img="labs/activity3-conditional/Activity3a-Conditional-Data-Table.png",cap="IF routes to Data Table or a thank-you email",
 steps=["Add an Attending? field (Yes/No dropdown) to the form.",
   "Add an IF node after the Form Trigger: condition {{ $json.Attending }} equals Yes.",
   "On the true output, add a Data Table → Insert Row to save the RSVP.",
   "On the false output, add a Gmail node sending a friendly thank-you.",
   "Save and keep the workflow Active."],
 test="Submit once with Attending = Yes (new Data Table row) and once with No (thank-you email)."))
activity_block(dict(tag="ACT 3b",title="Activity 3b — Conditional Response (Google Sheets / Excel)",kicker=K1,
 desc="Make data persistent by replacing the Data Table with Google Sheets (or Excel). Same logic — the Yes branch appends a row to a real spreadsheet you keep.",
 build="Form Trigger  →  IF  →  Google Sheets / Gmail",nodes="formTrigger, if, googleSheets, gmail",
 img="labs/activity3-conditional/Activity3b-Conditional-Google-Sheets.png",cap="IF routes to Google Sheets or a thank-you email",
 steps=["Create a Google Sheet 'Event RSVPs' with a header row: Name, Email, Phone, Date.",
   "Add a Google Sheets credential (OAuth2) in n8n and authorise it.",
   "Take Activity 3a; on the true branch replace the Data Table with a Google Sheets → Append Row node.",
   "Select your spreadsheet and sheet, then map each column to the form fields.",
   "Leave the false branch (thank-you email) unchanged; Save and keep Active."],
 test="Submit with Attending = Yes and confirm a new row is appended to your Google Sheet."))
brk("Lunch Break","1 hour  ·  see you at 2:00 pm",AMBER)

# ---------- TOPIC 2: AI AGENTS ----------
section("TOPIC 2","AI Agents and RAG","02","LLM · Memory · Tools · System Instruction")
content("What is Agentic AI?",[
 "Traditional automation follows fixed rules you wire by hand.",
 "An AI agent uses an LLM to understand language and decide what to do.",
 "It can call tools, remember context, and handle open-ended requests.",
 "Agentic = the AI takes actions toward a goal, not just answers once."],kicker="CONCEPT")
cards3("Why Use an AI Agent?",[
 (BLUE,"Understands language",["Handles free-text requests","No rigid forms or keywords","Summarises & reasons"]),
 (TEAL,"Takes action",["Calls tools & APIs","Looks up your data (RAG)","Chains multiple steps"]),
 (VIOLET,"Flexible",["One agent, many questions","Adapts to new cases","Easy to extend with tools"])],kicker="WHY AGENTS")
two_col("Popular LLM Models",[
 ("OpenAI - gpt-4.1, gpt-4.1-mini, gpt-4o",0),
 ("Google - Gemini 2.x Flash / Pro",0),
 ("Anthropic - Claude (Sonnet, Opus, Haiku)",0)],
 [("Meta - Llama 3.x (open weights)",0),
 ("Mistral - Mistral / Mixtral",0),
 ("This course uses OpenAI gpt-4.1-mini",0)],
 kicker="THE BRAINS",lhead="Hosted models",rhead="Open & course default")
content("AI Agent Using an LLM",[
 "The LLM is the 'brain' that interprets the user and plans a response.",
 "Tools give the agent abilities (look up data, search, call an API).",
 "Memory lets it hold a conversation across messages.",
 "A system instruction sets its role, tone and rules."],kicker="HOW IT WORKS")
two_col("AI Agent Concepts in n8n",[
 ("LLM — generates the replies",0),("OpenAI gpt-4.1-mini or Google Gemini",1),
 ("Memory — recalls the conversation",0),("Simple / window buffer memory",1)],
 [("Tools — actions the agent can call",0),("Data Table, Vector Store, HTTP, sub-workflow",1),
 ("System Instruction — role & rules",0),("Persona, scope, and do/don't guidance",1)],
 kicker="THE PIECES",lhead="Model & Memory",rhead="Tools & Instruction")
content("Creating System Prompts",[
 "State the agent's role: 'You are an HR admin assistant.'",
 "Give scope and limits: what it should and shouldn't answer.",
 "Tell it which tool to use for which kind of question.",
 "Keep it concise, specific and testable."],kicker="INSTRUCTION DESIGN")
content("Agent Instruction Best Practices",[
 "Be explicit about routing between multiple tools/sources.",
 "Ask for a clear format when you need structured output.",
 "Tell it to say 'I don't know' rather than invent answers.",
 "Iterate: test with real questions and refine the prompt."],kicker="BEST PRACTICE")
content("Memory & Tools",[
 "Add memory so follow-up questions keep context.",
 "Attach tools to the agent's Tool input; describe each tool clearly.",
 "The agent decides when to call a tool based on your description.",
 "All agents in this course use OpenAI gpt-4.1-mini."],kicker="EXTENDING THE AGENT")
K2="TOPIC 2 · AI AGENTS"
activity_block(dict(tag="ACT 4a",title="Activity 4a — Telegram-Triggered AI Agent (Customer Service)",kicker=K2,
 desc="Build your first AI agent: a customer-service chatbot you talk to from Telegram, using an LLM, memory and a system instruction.",
 build="Telegram Trigger  →  AI Agent (+ model + memory)  →  Telegram reply",nodes="telegramTrigger, agent, lmChatOpenAi, memory, telegram",
 img="labs/activity4-telegram-agent/Activity4a-Telegram-Agent.png",cap="Telegram-triggered AI agent with model + memory",
 steps=["In Telegram, open @BotFather → /newbot, then copy the bot token.",
   "Add a Telegram credential in n8n and paste the token.",
   "Add a Telegram Trigger node (fires on each message).",
   "Add an AI Agent node after the trigger.",
   "Attach an OpenAI Chat Model (gpt-4.1-mini) and a Simple Memory node.",
   "Write the system instruction (friendly customer-service assistant).",
   "Add a Telegram → Send Message: Chat ID = {{ $json.message.chat.id }}, Text = the agent output. Save & Activate."],
 test="Message your bot in Telegram and confirm it replies."))
activity_block(dict(tag="ACT 4b",title="Activity 4b — Telegram Agent + Data Table Tool (HR Admin)",kicker=K2,
 desc="Give the agent a tool: an employee Data Table it can query. The same bot now answers HR-admin questions from real data.",
 build="Telegram  →  AI Agent  +  Data Table Tool  →  reply",nodes="agent, dataTableTool, telegram",
 img="labs/activity4-telegram-agent/Activity4b-Telegram-Data-Table.png",cap="Agent with a Data Table tool",
 steps=["Create a Data Table 'Employees' and upload the provided employees.csv (100 records).",
   "Open your Activity 4a workflow.",
   "Add a Data Table Tool and attach it to the AI Agent's Tool input.",
   "Point the tool at the Employees table and describe it (look up staff by name/department).",
   "Update the system instruction to use the Employees tool for staff questions. Save & Activate."],
 test="Ask 'Which department is <a name from the CSV> in?' and confirm it answers from the table."))
content("End of Day 1 — Recap",[
 "You built form automations, conditional logic and data storage.",
 "You created a Telegram AI agent with memory and a Data Table tool.",
 "Tomorrow: RAG, Webhooks and external APIs."],kicker="WRAP-UP")

# ---------- TOPIC 3: RAG ----------
section("TOPIC 2 (cont.)","Retrieval-Augmented Generation (RAG)","02","Tokenization · Embeddings · Vector Stores")
content("What is RAG?",[
 "RAG lets an agent answer from YOUR documents, not just its training data.",
 "Documents are split, embedded and stored; relevant chunks are retrieved per question.",
 "Reduces hallucination and keeps answers grounded and current."],kicker="CONCEPT")
content("Text Embedding",[
 "Tokenization splits text into tokens the model can process.",
 "An embedding turns a chunk of text into a vector (list of numbers).",
 "Similar meanings produce vectors that are close together."],kicker="EMBEDDINGS")
content("Vector Database",[
 "Vectors are stored in a vector store (in-memory, Pinecone, etc.).",
 "At query time, the question is embedded and the closest chunks are retrieved.",
 "Those chunks are given to the LLM as context to answer."],kicker="VECTOR STORE")
img_slide("How RAG Works",IMG("courseware/assets/rag-flow.png"),
          "User → Prompt → Data Retrieval (search/retrieve over your data sources) → Generator → Response",
          kicker="TOPIC 2 · RAG")
K3="TOPIC 2 · RAG"
activity_block(dict(tag="ACT 5",title="Activity 5 — Add RAG to the Telegram Agent (Two Knowledge Sources)",kicker=K3,
 desc="Upgrade the agent to answer from documents (policies/FAQs) AND the Data Table. It must route to the right source for each question — two knowledge sources, one agent.",
 build="Telegram  →  AI Agent  +  Vector Store (RAG) + Data Table  →  reply",nodes="agent, vectorStoreInMemory, embeddingsOpenAi, dataTableTool",
 img="labs/activity5-rag/Activity5-RAG-Telegram.png",cap="Agent routes between a RAG vector store and a Data Table",
 steps=["Prepare documents: use MyCompany-HR-SOP.docx / IT-Support-FAQ.docx, or generate FAQs with Claude Code.",
   "Build the ingestion path: upload → Embeddings (OpenAI) → Vector Store (Insert) with a Default Data Loader.",
   "In your Telegram agent, add a Vector Store retrieval tool alongside the Data Table tool.",
   "Rewrite the system instruction to route: knowledge base for policy/FAQ, Data Table for staff records.",
   "Save and keep Active; have a few learners present their chatbot."],
 test="Ask a policy question and a staff-record question; confirm each is answered from the correct source."))
content("From In-Memory to a Real Vector Database",[
 "The in-memory store resets when the workflow restarts - fine for a demo.",
 "For production, use a hosted vector database that persists your embeddings.",
 "Pinecone is a popular managed vector database that scales to millions of vectors.",
 "Same idea: embed your documents once, then retrieve the closest chunks per question."],kicker="WHY PINECONE")
website_slide("Pinecone",IMG("courseware/assets/site-pinecone.png"),
 ["Pinecone is a managed (cloud) vector database for RAG.","Free 'Starter' tier is enough for this lab.","Create an index, then point n8n's Pinecone node at it."],
 kicker="VECTOR DATABASE")
content("Create a Pinecone Index",[
 "1. Sign up at pinecone.io and open the console.",
 "2. Create an API key (Database -> API Keys) - you'll paste it into n8n.",
 "3. Create an Index: give it a name (e.g. n8n-course).",
 "4. Set Dimensions to match your embeddings - OpenAI text-embedding-3-small = 1536.",
 "5. Use metric 'cosine'. The index name + key go into the n8n Pinecone node."],kicker="SETUP")
activity_overview("ACT 5b","Activity 5b — RAG with Pinecone (Persistent Vector Database)",
 "Swap the in-memory vector store for Pinecone so your knowledge base persists. Upload documents into a Pinecone index, then let the Telegram agent answer from it via a Vector Store tool.",
 "Upload -> Embeddings (OpenAI) -> Pinecone (insert)   |   Telegram -> AI Agent + Pinecone tool -> reply",
 "vectorStorePinecone, embeddingsOpenAi, toolVectorStore, agent",kicker="TOPIC 2 (cont.) - RAG")
img_slide("Activity 5b - Pinecone RAG Workflow",IMG("labs/activity5-rag/Activity5b-Pinecone-RAG.png"),
 "Telegram agent answering from a Pinecone vector store (gpt-4.1-mini)",kicker="TOPIC 2 (cont.) - RAG")
for i,t in enumerate([
 "Sign up at pinecone.io (free Starter tier) and open the console at app.pinecone.io.",
 "Go to API Keys → create / copy an API key — paste it into n8n later.",
 "Open Indexes → Create index; name it n8n-course.",
 "Set Dimensions = 1536 and Metric = cosine, then create the index.",
 "Import Activity5b-Pinecone-Upload.json into n8n.",
 "Add a Pinecone credential and select your n8n-course index on the Pinecone node.",
 "Add your OpenAI credential on the Embeddings node; provide documents and run to insert into Pinecone.",
 "Import Activity5b-Pinecone-RAG.json (Telegram → AI Agent + Pinecone tool → reply).",
 "Select the same Pinecone index and credential, your OpenAI (gpt-4.1-mini) and Telegram credentials.",
 "Save and toggle Active."],1):
    step_slide("TOPIC 2 (cont.) - RAG","Activity 5b — RAG with Pinecone (Persistent Vector Database)",i,10,t)
test_slide("Activity 5b — RAG with Pinecone (Persistent Vector Database)","Upload a document, then ask the Telegram bot a question only answerable from it - the answer is retrieved from Pinecone.","TOPIC 2 (cont.) - RAG")
brk("Lunch Break","1 hour",AMBER)

# ---------- TOPIC 4: WEBHOOKS ----------
section("TOPIC 3","Webhooks","03","External triggers for your workflows")
content("What is a Webhook?",[
 "A Webhook is a URL that external systems call to trigger your workflow.",
 "Use cases: website chat, form submissions, payments, app notifications.",
 "Pair the Webhook trigger with a Respond to Webhook node to reply."],kicker="CONCEPT")
content("How a Webhook Works",[
 "1. You activate an n8n Webhook node - it gives you a unique URL.",
 "2. An external system (website, app, Telegram) sends an HTTP request to that URL.",
 "3. n8n runs your workflow with the request data as the trigger input.",
 "4. A Respond to Webhook node sends a reply back to the caller.",
 "It is 'reverse' of an API call: the outside world calls YOU when an event happens."],kicker="EVENT -> WORKFLOW")
two_col("Webhook: GET vs POST",[
 ("GET",0),("Data comes in the URL query string",1),
 ("?name=Alice&topic=AI",1),("Good for simple links / browser visits",1)],
 [("POST",0),("Data comes in the request body (JSON)",1),
 ('{ "message": "hello" }',1),("Used by website chat & app callbacks",1)],
 kicker="TWO WAYS TO SEND DATA",lhead="GET",rhead="POST")
two_col("HTTP Request vs Webhook",[
 ("HTTP Request node",0),("YOUR workflow CALLS an external service",1),("You pull data when you choose",1)],
 [("Webhook node",0),("An external service CALLS your workflow",1),("You react when an event happens",1)],
 kicker="DIRECTION OF THE CALL",lhead="Outbound",rhead="Inbound")
content("Webhook Node & URL",[
 "Each Webhook node has a Test URL and a Production URL.",
 "Set the HTTP method (GET/POST) and a path.",
 "Set Allowed Origins (CORS) to * so a browser page can call it."],kicker="SETUP")
content("Webhook Authentication",[
 "None — open endpoint (fine for demos).",
 "Basic Auth — username & password.",
 "Header Auth — a secret key in a request header.",
 "JWT — signed tokens for stronger security."],kicker="SECURING WEBHOOKS")
K4="TOPIC 3 · WEBHOOKS"
activity_block(dict(tag="ACT 6",title="Activity 6 — Website Chatbot via Webhook (Investment Advisor)",kicker=K4,
 desc="Expose an AI agent to a public website via a Webhook. The Investment Advisor page has an enquiry form and a floating chatbot; both post to one n8n webhook.",
 build="Webhook  →  AI Agent  →  Respond to Webhook  (+ email the advisor)",nodes="webhook, agent, respondToWebhook, gmail",
 img="labs/activity6-investment-advisor/Activity6-Investment-Advisor.png",cap="One webhook, two paths: enquiry email + AI chat",
 steps=["Import Activity6-Investment-Advisor.json into n8n.",
   "Open the Webhook node(s) and set Allowed Origins (CORS) to *.",
   "Re-select your OpenAI and Gmail credentials on the agent and email nodes.",
   "Review the compliance system instruction (no guaranteed returns, no personalised advice).",
   "Save, Activate, and copy the webhook Production URL.",
   "Open index.html, click the gear, and paste your webhook URL; have learners present their site."],
 test="On the website, send a chat message and submit the enquiry form; confirm the bot replies and the advisor gets the email."))
brk("Tea Break","15 minutes",TEAL)

# ---------- TOPIC 5: API ----------
section("TOPIC 4","API and HTTP Request","04","Pull live data from external services")
content("What is an API?",[
 "An API lets your workflow request data from another service over HTTP.",
 "You send a request; the service sends back a response (usually JSON).",
 "APIs power live data: prices, weather, news, CRM records."],kicker="CONCEPT")
content("How an API Works",[
 "1. Your workflow (the client) sends an HTTP request to an API endpoint (URL).",
 "2. You include a method, headers (often an API key) and parameters.",
 "3. The server processes it and returns a response - usually JSON.",
 "4. n8n parses the JSON and passes the fields to the next node.",
 "Here, YOU call the outside service (the opposite of a webhook)."],kicker="WORKFLOW -> SERVICE")
two_col("API: GET vs POST",[
 ("GET - read data",0),("Fetch prices, news, records",1),
 ("Parameters in the URL query",1),("Safe & repeatable",1)],
 [("POST - send / create data",0),("Submit a form, create a record",1),
 ("Parameters in the JSON body",1),("Changes data on the server",1)],
 kicker="THE TWO MAIN METHODS",lhead="GET",rhead="POST")
content("Components of an HTTP Request",[
 "Method — GET (read), POST (send), PUT, DELETE.",
 "URL / endpoint — the address of the resource.",
 "Headers — metadata, including authentication.",
 "Query params / body — the inputs you send."],kicker="ANATOMY")
content("HTTP Response & Status Codes",[
 "Body — the data returned (often JSON).",
 "200 OK — success; 201 Created.",
 "401 Unauthorized — bad/missing API key; 429 — rate limited.",
 "4xx = your request; 5xx = the server."],kicker="READING RESPONSES")
content("HTTP Request Node",[
 "Configure method, URL, headers and query parameters.",
 "Store API keys in credentials, never hard-coded.",
 "Parse the JSON response and pass fields to the next node."],kicker="IN n8n")
K5="TOPIC 4 · API & HTTP"
website_slide("Twelve Data",IMG("courseware/assets/site-twelvedata.png"),
 ["Twelve Data provides live stock / forex / crypto market data.","Sign up (free Basic plan), then Account -> API Keys.","Paste the apikey into the 3 'candles' HTTP nodes."],kicker="MARKET DATA API")
website_slide("NewsAPI",IMG("courseware/assets/site-newsapi.png"),
 ["NewsAPI returns recent news articles for a search query.","Register (free Developer plan) and copy your API key.","Store it as a Query Auth credential (name = apiKey) on the news node."],kicker="NEWS API")
activity_block(dict(tag="ACT 7",title="Activity 7 — Finance API → Telegram (AI Day-Trading Agent)",kicker=K5,
 desc="Ask the Telegram bot about a stock; it resolves the ticker, pulls candles from Twelve Data and headlines from NewsAPI, and replies with a Buy/Sell/Hold call and reasoning.",
 build="Telegram → Extract ticker → HTTP (Twelve Data + NewsAPI) → AI Agent → reply",nodes="telegram, httpRequest, agent, lmChatOpenAi",
 img="labs/activity7-finance-advisor/Activity7-Finance-Advisor.png",cap="Multi-timeframe candles + news → AI day-trading agent",
 steps=["Sign up at twelvedata.com (free Basic plan) → Account → API Keys → copy your key.",
   "Sign up at newsapi.org (free Developer plan) → copy your key from newsapi.org/account.",
   "Import Activity7-Finance-Advisor.json into n8n.",
   "Open candles1min → Query Parameters → find apikey → replace with your Twelve Data key.",
   "Repeat for candles15min and candles1hr — all three nodes need the same Twelve Data key.",
   "Open the news node → Authentication: Generic Credential Type → Query Auth.",
   "Credential → Create New; set Name = apiKey and Value = your NewsAPI key, then Save.",
   "Re-select your OpenAI and Telegram credentials on the model and Telegram nodes.",
   "Review: Telegram Trigger → Extract Ticker → candles (1m/15m/1h) + news → AI Agent → reply.",
   "Save and toggle Active; optionally open index.html and paste your Twelve Data key + bot username."],
 test="Message the bot 'Should I buy AAPL?' and confirm it returns a recommendation with reasoning."))
content("End of Day 2 — Recap",[
 "You grounded an agent with RAG over your own documents.",
 "You exposed an agent to a website via a webhook.",
 "You pulled live market + news data through APIs into a Telegram agent."],kicker="WRAP-UP")

# ---------- TOPIC 6: SECURITY ----------
section("TOPIC 5","Security and Guardrails","05","Human-in-the-loop · Pre/Post guardrails")
content("Human in the Loop",[
 "Some actions are too sensitive to fully automate — money, hiring, sending on behalf.",
 "A human-in-the-loop step pauses the workflow for a person to Approve or Reject.",
 "n8n's Send and Wait for Response captures that decision (email / Telegram)."],kicker="CONCEPT")
K6="TOPIC 5 · SECURITY"
activity_block(dict(tag="ACT 8a",title="Activity 8a — Human-in-the-Loop Approval (Leave Application)",kicker=K6,
 desc="Model a leave-application approval: a request comes in, a manager is asked to approve, and the flow only continues — recording the leave and confirming — on approval.",
 build="Form → Manager Approval (Send & Wait) → IF → confirm / decline",nodes="formTrigger, gmail (sendAndWait), if",
 img="labs/activity8-guardrails/Activity8a-Human-in-the-Loop.png",cap="Approval pauses the flow until a manager decides",
 steps=["Start with a Form Trigger collecting Employee, Email, Dates, Reason.",
   "Add a Gmail → Send and Wait for Response (Approval) addressed to the manager.",
   "Add an IF node on the approval result.",
   "On Approved: record the leave (Data Table) and email a confirmation.",
   "On Rejected: email the employee that the request was declined. Save & Activate."],
 test="Submit a leave request, approve it from the manager email, and confirm the employee gets a confirmation."))
content("Guardrails",[
 "Pre-guardrail — validate/sanitise the input (block prompt-injection, PII) before the LLM.",
 "Post-guardrail — check the output (no secrets, no disallowed content) before sending.",
 "On a violation, route to a safe fallback or human review."],kicker="CONCEPT")
activity_block(dict(tag="ACT 8b",title="Activity 8b — Pre & Post Guardrails for the AI Agent",kicker=K6,
 desc="Wrap an AI agent so unsafe input never reaches the model and unsafe output never reaches the user. Add a pre-check before the agent and a post-check after it.",
 build="Webhook → Pre-check → AI Agent → Post-check → Respond / Blocked",nodes="webhook, if, agent, respondToWebhook",
 img="labs/activity8-guardrails/Activity8b-Guardrails.png",cap="Pre/post checks gate the agent",
 steps=["Start from the Activity 6 webhook agent (or the Telegram agent).",
   "Before the AI Agent, add a Guardrails node that checks the user message for secret keys or blocked keywords.",
   "After the AI Agent, add a second Guardrails check on the reply; adjust the keyword list as you like.",
   "If either guardrail fails, the false branch returns a safe canned response.",
   "Only send the reply when both guardrails pass.",
   "Save and keep Active."],
 test="Send a normal question (passes through) and a disallowed one — e.g. 'my password is sk-12345678' — and confirm the pre-guardrail blocks it with a safe reply."))
brk("Lunch Break","1 hour",AMBER)

# ---------- TOPIC 7: CAPSTONE ----------
section("TOPIC 6","Mini Capstone Project","06","Design · Build · Present · Assess")
content("Mini Capstone Project",[
 "In small groups, design and build an end-to-end automation using what you learned.",
 "Include: a trigger, an AI agent with a tool or RAG source, an external API or storage, and a guardrail / human-in-the-loop step.",
 "A worked example — Issue Reporting (form + image → Postgres + gallery) — is provided.",
 "Deliverables: a working workflow, a short demo, and a 3–5 minute presentation."],kicker="BRING IT TOGETHER")
content("Presentation & Assessment",[
 "Present the problem, your design, and what you'd improve.",
 "Written Assessment (SAQ) — 1 hour · Practical Performance (PP) — 1 hour.",
 "Open book: slides, Learner Guide and approved materials.",
 "Remember to take the Assessment digital attendance (TRAQOM)."],kicker="WRAP-UP")

# ---------- CLOSING ----------
s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
rect(s,Inches(0.85),Inches(2.7),Inches(0.14),Inches(1.5),TEAL)
txt(s,Inches(1.25),Inches(2.7),Inches(11),Inches(1.2),[[("Thank You!",48,INK,True)]])
txt(s,Inches(1.27),Inches(4.1),Inches(11.5),Inches(1.4),
 [[("Keep your local n8n running and keep building your own agents.",18,GREY,False)],
  [("Download all flows: github.com/tertiarycourses/TGS-2023035977-Agentic-AI-Automation-with-n8n",13,BLUE,True)],
  [("Powered by Tertiary Infotech Academy Pte Ltd  ·  www.tertiarycourses.com.sg",12,GREY,False)]],space=8)
footer(s)

OUT="courseware/Agentic AI Automation with n8n.pptx"
prs.save(f"{REPO}/{OUT}")
print("Saved:",OUT,"| slides:",len(prs.slides._sldIdLst))
